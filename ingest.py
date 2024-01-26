"""Load text from files,  split, ingest into vectorstore.
"""
import os
from fastapi import HTTPException
import requests

from pdfminer.high_level import extract_text_to_fp
from pdfminer.layout import LAParams
import io
from langchain.embeddings import OpenAIEmbeddings
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.vectorstores.faiss import FAISS
from langchain.docstore.document import Document
from langchain.document_loaders import UnstructuredHTMLLoader
from langchain.document_loaders.csv_loader import CSVLoader
from url_csv_loader import UrlCSVLoader
from langchain.document_loaders import UnstructuredPowerPointLoader
from langchain.document_loaders import UnstructuredWordDocumentLoader
from langchain.document_loaders import YoutubeLoader
import requests
from urllib.parse import urlparse
from pptx import Presentation
import mammoth
from bs4 import BeautifulSoup
import urllib.request
import re
import tiktoken
from langchain.text_splitter import RecursiveCharacterTextSplitter
from urllib.parse import urlparse, parse_qs, urlencode, urlunparse
import unicodedata
import pickle
from langchain.vectorstores.base import VectorStore
from pathlib import Path
from langchain.vectorstores import Pinecone
import pinecone
from typing import List
import assemblyai as aai





def normalize_string(input_string): # remove non-ascii and other characters from namespace name for Pinecone Standard
    normalized_string = unicodedata.normalize('NFKD', input_string)
    s= ''.join([char for char in normalized_string if not unicodedata.combining(char)])
    return ''.join([char for char in s if (ord(char) < 128) and (ord(char)>=32)])


def num_tokens_from_string(string: str, model:str) -> int:
    """Returns the number of tokens in a text string assuming a given GPT Model using tiktoken module"""
    model=model.lower()
    if "gpt-4" in model:
        model="gpt-4"
    encoding=tiktoken.encoding_for_model(model)
    num_tokens = len(encoding.encode(string))
    return num_tokens
        
def isvalid_url(url: str) -> bool:
        """Check if the url is valid."""
        parsed = urlparse(url)
        return bool(parsed.netloc) and bool(parsed.scheme)
    
def extract_video_id(url):
    """
    The extract_video_id function is a Python function that extracts the video ID from various YouTube URL formats. 
    It utilizes regular expressions to match different URL patterns commonly used by YouTube, including standard videos, shorts, live streams, and more.

Parameters
url (str): The YouTube URL from which you want to extract the video ID.

Return Value
The function returns the extracted video ID as a string. If no video ID is found in the URL, it returns None.
    """
    
    # Regular expression patterns to match different URL formats
    patterns = [
        r"(?:youtu\.be\/|youtube\.com\/(?:embed\/|v\/|watch\?v=|watch\?.+&v=))([\w-]+)",
        r"(?:youtube\.com\/shorts\/|youtu\.be\/shorts\/)([\w-]+)",
        r"(?:youtube\.com\/live\/|youtu\.be\/live\/)([\w-]+)",
        r"(?:youtube\.com\/watch\?.*?v=|youtu\.be\?)([\w-]+)"
    ]
    
    for pattern in patterns:
        match = re.search(pattern, url)
        if match:
            return match.group(1)

    # Return None if no video ID is found
    return None
 
def is_youtube_video(url): # check if the url is a youtube video url
    videoId=extract_video_id(url=url)
    if videoId!=None :
        response = requests.head(url)
    else:
        response=None
    return videoId!=None and response!=None and response.status_code == 200

def ingest_doc_to_local_vectstore(doc:str,chunkSize, chunkOverlap, open_ai_key):
        filename=doc.strip()
        filename=filename.replace('https:https:', 'https:')
        filename=filename.replace('https://https:', 'https:')
        extension = os.path.splitext(filename)[1].lower()
        if extension==".pdf" or filename[-3:].lower()=="pdf":
                raw_documents=pdf_to_rawdocs(filename)
        elif extension==".csv" or filename[-3:].lower()=="csv":
            raw_documents=csv_to_rawdocs(filename)
        elif extension=='.ppt' or extension=='.pptx' or filename[-3:].lower()=="ppt" or filename[-4:].lower()=="pptx":
            raw_documents=pptx_to_rawdocs(filename)
        elif extension=='.doc' or extension=='.docx' or filename[-3:].lower()=="doc" or filename[-4:].lower()=="docx":
            raw_documents=docx_to_rawdocs(filename)
        elif extension=='.html' or extension=='.htm' or filename[-4:].lower()=="html" or filename[-3:].lower()=="htm":
            raw_documents=html_to_rawdocs(filename)
        elif extension=='.txt'or filename[-3:].lower()=="txt":
            raw_documents=txt_to_rawdocs(filename)
        elif is_youtube_video(filename):
            raw_documents=youtube_to_rawdocs(filename)
        else:
            raise HTTPException(status_code=400, detail="File path \'" +filename+"\' is not a valid file or url, or format not yet supported")
    
        text_splitter = RecursiveCharacterTextSplitter(chunk_size=chunkSize, chunk_overlap=chunkOverlap,)
        documents = text_splitter.split_documents(raw_documents)
        embeddings = OpenAIEmbeddings(openai_api_key=open_ai_key)
        vectorstore = FAISS.from_documents(documents, embeddings)

        # Save vectorstore
        with open("vectorstore.pkl", "wb") as f:
            pickle.dump(vectorstore, f)
        return vectorstore

def ingest_folder_to_local_vectstore(folder:str,chunkSize, chunkOverlap, open_ai_key):
    documents = []
    not_ingested=[]
    for doc in os.listdir(folder): 
        filename=folder+"/"+doc
        filename=filename.strip()
        filename=filename.replace('https:https:', 'https:')
        filename=filename.replace('https://https:', 'https:')
        extension = os.path.splitext(filename)[1].lower()
        if extension==".pdf" or filename[-3:].lower()=="pdf":
                raw_documents=pdf_to_rawdocs(filename)
        elif extension==".csv" or filename[-3:].lower()=="csv":
            raw_documents=csv_to_rawdocs(filename)
        elif extension=='.ppt' or extension=='.pptx' or filename[-3:].lower()=="ppt" or filename[-4:].lower()=="pptx":
            raw_documents=pptx_to_rawdocs(filename)
        elif extension=='.doc' or extension=='.docx' or filename[-3:].lower()=="doc" or filename[-4:].lower()=="docx":
            raw_documents=docx_to_rawdocs(filename)
        elif extension=='.html' or extension=='.htm' or filename[-4:].lower()=="html" or filename[-3:].lower()=="htm":
            raw_documents=html_to_rawdocs(filename)
        elif extension=='.txt'or filename[-3:].lower()=="txt":
            raw_documents=txt_to_rawdocs(filename)
        elif is_youtube_video(filename):
            raw_documents=youtube_to_rawdocs(filename)
        else:
            not_ingested.append(doc)
            continue
        documents.extend(raw_documents)
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=chunkSize, chunk_overlap=chunkOverlap,)
    documents = text_splitter.split_documents(documents)
    embeddings = OpenAIEmbeddings(openai_api_key=open_ai_key)
    vectorstore = FAISS.from_documents(documents, embeddings)
    # Save vectorstore
    with open("vectorstore.pkl", "wb") as f:
            pickle.dump(vectorstore, f)
    return vectorstore, not_ingested
    
def add_doc_to_local_vectstore(filename:str):
    if not Path("vectorstore.pkl").exists():
        print("vectorstore.pkl does not exist, please run Ingest first")
        return None, False
    if not Path(filename).exists():
        print(filename+" does not exist (you may have to input the full path)")
        return None, False
    vectstore:VectorStore
    with open("vectorstore.pkl", "rb") as f:
        vectstore = pickle.load(f)
    filename=filename.strip()
    filename=filename.replace('https:https:', 'https:')
    filename=filename.replace('https://https:', 'https:')
    extension = os.path.splitext(filename)[1].lower()
    if extension==".pdf" or filename[-3:].lower()=="pdf":
                raw_documents=pdf_to_rawdocs(filename)
    elif extension==".csv" or filename[-3:].lower()=="csv":
            raw_documents=csv_to_rawdocs(filename)
    elif extension=='.ppt' or extension=='.pptx' or filename[-3:].lower()=="ppt" or filename[-4:].lower()=="pptx":
            raw_documents=pptx_to_rawdocs(filename)
    elif extension=='.doc' or extension=='.docx' or filename[-3:].lower()=="doc" or filename[-4:].lower()=="docx":
            raw_documents=docx_to_rawdocs(filename)
    elif extension=='.html' or extension=='.htm' or filename[-4:].lower()=="html" or filename[-3:].lower()=="htm":
            raw_documents=html_to_rawdocs(filename)
    elif extension=='.txt'or filename[-3:].lower()=="txt":
            raw_documents=txt_to_rawdocs(filename)
    elif is_youtube_video(filename):
            raw_documents=youtube_to_rawdocs(filename)
    else:
            return None, False
    vectstore.add_documents(raw_documents) # to add documents to an existing vectorstore
    # Save vectorstore
    with open("vectorstore.pkl", "wb") as f:
            pickle.dump(vectstore, f)
    return vectstore, True

def ingest_folder_to_pinecone(folder:str,chunkSize, chunkOverlap, ind_name, nsname, delete_ns_if_exists, openaikey, pineconekey,pineconeenv):
    documents = []
    not_ingested=[]
    for doc in os.listdir(folder): 
        filename=folder+"/"+doc
        filename=filename.strip()
        filename=filename.replace('https:https:', 'https:')
        filename=filename.replace('https://https:', 'https:')
        extension = os.path.splitext(filename)[1].lower()
        if extension==".pdf" or filename[-3:].lower()=="pdf":
                raw_documents=pdf_to_rawdocs(filename)
        elif extension==".csv" or filename[-3:].lower()=="csv":
            raw_documents=csv_to_rawdocs(filename)
        elif extension=='.ppt' or extension=='.pptx' or filename[-3:].lower()=="ppt" or filename[-4:].lower()=="pptx":
            raw_documents=pptx_to_rawdocs(filename)
        elif extension=='.doc' or extension=='.docx' or filename[-3:].lower()=="doc" or filename[-4:].lower()=="docx":
            raw_documents=docx_to_rawdocs(filename)
        elif extension=='.html' or extension=='.htm' or filename[-4:].lower()=="html" or filename[-3:].lower()=="htm":
            raw_documents=html_to_rawdocs(filename)
        elif extension=='.txt'or filename[-3:].lower()=="txt":
            raw_documents=txt_to_rawdocs(filename)
        elif is_youtube_video(filename):
            raw_documents=youtube_to_rawdocs(filename)
        else:
            not_ingested.append(doc)
            continue
        documents.extend(raw_documents)
        text_splitter = RecursiveCharacterTextSplitter(chunk_size=chunkSize, chunk_overlap=chunkOverlap,)
        documents = text_splitter.split_documents(documents)
        # embeddings = OpenAIEmbeddings(openai_api_key=os.environ.get('OPENAI_API_KEY'))
        
        embeddings = OpenAIEmbeddings(openai_api_key=openaikey)
        # initialize pinecone
        pinecone.init(
            api_key=pineconekey,  # find at app.pinecone.io
            environment=pineconeenv,  
            namespace=nsname)
    if (delete_ns_if_exists):
        pinecone.Index(index_name=ind_name).delete(delete_all=True, namespace=nsname)
    return Pinecone.from_documents(documents, embeddings, index_name=ind_name, namespace=nsname), not_ingested

def ingest_urls_and_text_to_pinecone(urls:List[str],chunkSize, chunkOverlap, ind_name, nsname, delete_ns_if_exists, openaikey, pineconekey,pineconeenv,
                                     text=None):
    documents = []
    if urls!=None:
     for filename in urls: 
        filename=filename.strip()
        filename=filename.replace('https:https:', 'https:')
        filename=filename.replace('https://https:', 'https:')
        extension = os.path.splitext(filename)[1].lower()
        if extension==".pdf" or filename[-3:].lower()=="pdf":
                raw_documents=pdf_to_rawdocs(filename)
        elif extension==".csv" or filename[-3:].lower()=="csv":
            raw_documents=csv_to_rawdocs(filename)
        elif extension=='.ppt' or extension=='.pptx' or filename[-3:].lower()=="ppt" or filename[-4:].lower()=="pptx":
            raw_documents=pptx_to_rawdocs(filename)
        elif extension=='.doc' or extension=='.docx' or filename[-3:].lower()=="doc" or filename[-4:].lower()=="docx":
            raw_documents=docx_to_rawdocs(filename)
        elif extension=='.html' or extension=='.htm' or filename[-4:].lower()=="html" or filename[-3:].lower()=="htm":
            raw_documents=html_to_rawdocs(filename)
        elif extension=='.txt'or filename[-3:].lower()=="txt":
            raw_documents=txt_to_rawdocs(filename)
        elif extension=='.mp3'or filename[-3:].lower()=="mp3":
            raw_documents=mp3_to_rawdocs(filename)
        elif is_youtube_video(filename):
            raw_documents=youtube_to_rawdocs(filename)
        else:
            continue
        documents.extend(raw_documents)
    if (text!=None and text!=""):
        raw_documents=string_to_rawdocs(text,None)
        documents.extend(raw_documents)
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=chunkSize, chunk_overlap=chunkOverlap,)
    documents = text_splitter.split_documents(documents)        
    embeddings = OpenAIEmbeddings(openai_api_key=openaikey)
    # initialize pinecone
    pinecone.init(
            api_key=pineconekey,  # find at app.pinecone.io
            environment=pineconeenv,  
            namespace=nsname)
    if (delete_ns_if_exists):
        pinecone.Index(index_name=ind_name).delete(delete_all=True, namespace=nsname)
    return Pinecone.from_documents(documents, embeddings, index_name=ind_name, namespace=nsname)

def add_doc_to_pinecone(filename:str, chunkSize, chunkOverlap, ind_name, nsname, openaikey, pineconekey,pineconeenv):
    if not Path(filename).exists():
        print(filename+" does not exist (you may have to input the full path)")
        return None, False
    filename=filename.strip()
    filename=filename.replace('https:https:', 'https:')
    filename=filename.replace('https://https:', 'https:')
    extension = os.path.splitext(filename)[1].lower()
    if extension==".pdf" or filename[-3:].lower()=="pdf":
                raw_documents=pdf_to_rawdocs(filename)
    elif extension==".csv" or filename[-3:].lower()=="csv":
            raw_documents=csv_to_rawdocs(filename)
    elif extension=='.ppt' or extension=='.pptx' or filename[-3:].lower()=="ppt" or filename[-4:].lower()=="pptx":
            raw_documents=pptx_to_rawdocs(filename)
    elif extension=='.doc' or extension=='.docx' or filename[-3:].lower()=="doc" or filename[-4:].lower()=="docx":
            raw_documents=docx_to_rawdocs(filename)
    elif extension=='.html' or extension=='.htm' or filename[-4:].lower()=="html" or filename[-3:].lower()=="htm":
            raw_documents=html_to_rawdocs(filename)
    elif extension=='.txt'or filename[-3:].lower()=="txt":
            raw_documents=txt_to_rawdocs(filename)
    elif is_youtube_video(filename):
            raw_documents=youtube_to_rawdocs(filename)
    else:
            return None, False
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=chunkSize, chunk_overlap=chunkOverlap,)
    documents = text_splitter.split_documents(raw_documents)
    embeddings = OpenAIEmbeddings(openai_api_key=openaikey)
    # initialize pinecone
    pinecone.init(
            api_key=pineconekey,  # find at app.pinecone.io
            environment=pineconeenv,  
            namespace=nsname)
    return Pinecone.from_documents(documents, embeddings, index_name=ind_name, namespace=nsname), True

def add_string_to_pinecone(text:str, chunkSize, chunkOverlap, ind_name, nsname, openaikey, pineconekey,pineconeenv):
    raw_documents=string_to_rawdocs(text)
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=chunkSize, chunk_overlap=chunkOverlap,)
    documents = text_splitter.split_documents(raw_documents)
    embeddings = OpenAIEmbeddings(openai_api_key=openaikey)
    # initialize pinecone
    pinecone.init(
            api_key=pineconekey,  # find at app.pinecone.io
            environment=pineconeenv,  
            namespace=nsname)
    Pinecone.from_documents(documents, embeddings, index_name=ind_name, namespace=nsname)
   
def pdf_to_rawdocs(filename, metadata=None):
    """ 
    The pdf_to_rawdocs function is a Python function that converts PDF files into a list of raw documents. 
    It can accept a local file path or a URL pointing to a PDF file. 
    The function extracts the text content from the PDF and wraps it in a format suitable for further processing and ingestion, 
    along with optional metadata.

Parameters
filename (str): The path to a local PDF file or a URL pointing to a PDF file.
metadata (dict, optional): A dictionary containing metadata information associated with the PDF document. If not provided, a default metadata dictionary will be created with the PDF source information.

Return Value
The function returns a list of raw documents, where each document is represented as an instance of the Document class. Each Document instance contains the extracted text content from a PDF page and associated metadata.
    """ 

    try:
        if not os.path.isfile(filename) and isvalid_url(filename):
        # Download the PDF from the URL
            response = requests.get(filename)
            pdf_file = io.BytesIO(response.content)
        elif os.path.isfile(filename) and not isvalid_url(filename):
            pdf_file = open(filename, "rb")
        else:
            raise HTTPException(status_code=400, detail="File path " +filename+" is not a valid file or url")
        # Extract the text from the PDF
        output_string = io.StringIO()
        with pdf_file as fp:
                extract_text_to_fp(
                    fp,  # type: ignore[arg-type]
                    output_string,
                    codec="",
                    laparams=LAParams(),
                    output_type="text",
                )
        if metadata==None:
            metadata={"source": filename}
        raw_documents= [Document(page_content=output_string.getvalue(), metadata=metadata)]
        return raw_documents
    except Exception:
            raise HTTPException(status_code=400, detail="File path " +filename+" is not a valid file or url")

def csv_to_rawdocs(filename, metadata=None):
        """
        The csv_to_rawdocs function is a Python function designed to convert CSV files into a list of raw documents. 
        The function can handle both local file paths and URLs pointing to CSV files. 
        It uses appropriate loaders based on the file type to extract data from the CSV and wraps it in a format suitable 
        for further processing and ingestion, along with optional metadata.

Parameters
filename (str): The path to a local CSV file or a URL pointing to a CSV file.
metadata (dict, optional): A dictionary containing metadata information associated with the CSV data. If not provided, a default metadata dictionary will be created with the CSV source information.

Return Value
The function returns a list of raw documents, where each document is represented as an instance of the Document class. Each Document instance contains the extracted content from a CSV row and associated metadata.
        """
        if metadata==None:
            metadata={"source": filename}
        if not os.path.isfile(filename) and isvalid_url(filename):
            loader = UrlCSVLoader(file_path=filename)
            raw_docs = loader.load()
        elif os.path.isfile(filename) and not isvalid_url(filename):
            loader = CSVLoader(file_path=filename)
            raw_docs = loader.load()
        elif not os.path.isfile(filename):
            raise HTTPException(status_code=400, detail="File path " +filename+" is not a valid file or url")
        if raw_docs!=None:
            raw_documents = [Document(page_content=raw_doc.page_content,metadata=metadata) for raw_doc in raw_docs]
        else:
            raise HTTPException(status_code=400, detail="File path " +filename+" generated no document")
        return raw_documents
    
def pptx_to_rawdocs(filename, metadata=None):  
    """
    The pptx_to_rawdocs function is a Python function that converts PowerPoint (PPTX) files into a list of raw documents. 
    The function can handle both local file paths and URLs pointing to PPTX files. 
    It uses appropriate techniques to extract text from the slides and wraps the extracted content in a format 
    suitable for further processing and ingestion, along with optional metadata.

Parameters
filename (str): The path to a local PPTX file or a URL pointing to a PPTX file.
metadata (dict, optional): A dictionary containing metadata information associated with the PPTX file. If not provided, a default metadata dictionary will be created with the PPTX source information.

Return Value
The function returns a list of raw documents, where each document represents the text content extracted from a slide of the PowerPoint presentation. Each document is represented as an instance of the Document class.
    """

    try:
        if metadata==None:
            metadata={"source": filename}  
        if not os.path.isfile(filename) and isvalid_url(filename):
            # Download the Powepoint from the URL
            response = requests.get(filename)
            pptx_file = io.BytesIO(response.content)
            prs = Presentation(pptx_file)
            # Extract text from all slides and put them into raw_documents
            documents = []
            for slide in prs.slides:
                slide_text = ''
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        slide_text += shape.text
                documents.append(slide_text)
            # Optionally, close the file-like object
            pptx_file.close()
            raw_documents = [Document(page_content=raw_doc,metadata=metadata) for raw_doc in documents]
        elif os.path.isfile(filename) and not isvalid_url(filename):
            loader = UnstructuredPowerPointLoader(filename)
            raw_docs = loader.load()
            raw_documents = [Document(page_content=raw_doc.page_content,metadata=metadata) for raw_doc in raw_docs]
        elif not os.path.isfile(filename):
            raise HTTPException(status_code=400, detail="File path" +filename+" is not a valid file or url") 
        return raw_documents
    except Exception:
            raise HTTPException(status_code=400, detail="File path " +filename+" is not a valid file or url")
    
def docx_to_rawdocs(filename, metadata=None):
    """
    The docx_to_rawdocs function is a Python function that converts Microsoft Word (DOCX) files into a list of raw documents. 
    The function supports both local file paths and URLs pointing to DOCX files. 
    It uses suitable techniques to extract text from the DOCX file and wraps the extracted content in a format 
    suitable for further processing and ingestion, along with optional metadata.

Parameters
filename (str): The path to a local DOCX file or a URL pointing to a DOCX file.
metadata (dict, optional): A dictionary containing metadata information associated with the DOCX file. If not provided, a default metadata dictionary will be created with the DOCX source information.

Return Value
The function returns a list of raw documents, where each document represents the text content extracted from the DOCX file. Each document is represented as an instance of the Document class.
    """
    try:
        if metadata==None:
            metadata={"source": filename}  
        if not os.path.isfile(filename) and isvalid_url(filename):
            # Download the Word document from the URL
            response = requests.get(filename)
            docx_file = io.BytesIO(response.content)
            result = mammoth.extract_raw_text(docx_file)
            documents = [result.value]
            raw_documents = [Document(page_content=raw_doc,metadata=metadata) for raw_doc in documents]
        elif os.path.isfile(filename) and not isvalid_url(filename):
            loader = UnstructuredWordDocumentLoader(filename)
            raw_docs = loader.load()
            raw_documents = [Document(page_content=raw_doc.page_content,metadata=metadata) for raw_doc in raw_docs]
        elif not os.path.isfile(filename):
            raise HTTPException(status_code=400, detail="File path " +filename+" is not a valid file or url") 
        return raw_documents
    except Exception:
            raise HTTPException(status_code=400, detail="File path " +filename+" is not a valid file or url")
    
def html_to_rawdocs(filename, metadata=None):
    """
    The html_to_rawdocs function is a Python function that converts HTML files or web pages into a list of raw documents. 
    The function can handle both local file paths and URLs pointing to HTML content. 
    It uses suitable techniques to extract plain text from the HTML and wraps the extracted content in a format 
    suitable for further processing and ingestion, along with optional metadata.

Parameters
filename (str): The path to a local HTML file or a URL pointing to an HTML page.
metadata (dict, optional): A dictionary containing metadata information associated with the HTML content. If not provided, a default metadata dictionary will be created with the HTML source information.

Return Value
The function returns a list of raw documents, where each document represents the plain text content extracted from the HTML content. Each document is represented as an instance of the Document class.
    """

    try:
        if metadata==None:
            metadata={"source": filename}  
        if not os.path.isfile(filename) and isvalid_url(filename):
            r=requests.get(url=filename)
            contentType=r.headers['Content-Type'].lower()
            if 'utf-8' in contentType.lower():
                soup = BeautifulSoup(r.content, 'lxml', from_encoding='utf-8') # or 'html.parser' but 'lxml' is faster. it also could be None: let BeautifulSoup choose the parser
            else:
                soup = BeautifulSoup(r.content, 'lxml')
            plain_text = soup.get_text()
            raw_documents = [Document(page_content=plain_text, metadata=metadata)]
        elif os.path.isfile(filename) and not isvalid_url(filename):
            loader = UnstructuredHTMLLoader(filename)
            raw_docs = loader.load()
            raw_documents = [Document(page_content=raw_doc.page_content,metadata=metadata) for raw_doc in raw_docs]
        elif not os.path.isfile(filename):
            raise HTTPException(status_code=400, detail="File path " +filename+" is not a valid file or url") 
        return raw_documents
    except Exception:
            raise HTTPException(status_code=400, detail="File path " +filename+" is not a valid file or url")
    
def txt_to_rawdocs(filename, metadata=None):
    """
    The txt_to_rawdocs function is a Python function that converts plain text files or textual content from URLs into a list of raw documents. 
    The function can handle both local file paths and URLs pointing to plain text files. 
    It reads the text content and wraps it in a format suitable for further processing and ingestion, along with optional metadata.

Parameters
filename (str): The path to a local plain text file or a URL pointing to a plain text file.
metadata (dict, optional): A dictionary containing metadata information associated with the text content. If not provided, a default metadata dictionary will be created with the source information.

Return Value
The function returns a list of raw documents, where each document represents the plain text content extracted from the text file or URL. Each document is represented as an instance of the Document class.
    """
    try:
        if metadata==None:
            metadata={"source": filename} 
        if not os.path.isfile(filename) and isvalid_url(filename):
            # Open the URL and read the contents of the file into a string
            with urllib.request.urlopen(filename) as response:
                text = response.read().decode('utf-8')
            documents = [text]
            raw_documents = [Document(page_content=raw_doc,metadata=metadata) for raw_doc in documents]
        elif os.path.isfile(filename) and not isvalid_url(filename):
            with open(filename, "r") as f:
                text = f.read()
            f.close()
            documents = [text]
            raw_documents = [Document(page_content=raw_doc,metadata=metadata) for raw_doc in documents]
       
        elif not os.path.isfile(filename):
            raise HTTPException(status_code=400, detail="File path " +filename+" is not a valid file or url") 
        return raw_documents
    except Exception:
            raise HTTPException(status_code=400, detail="File path " +filename+" is not a valid file or url")
    
def string_to_rawdocs(text, metadata=None):
    """
    The string_to_rawdocs function is a Python function that converts plain text string into a list of raw documents 
    suitable for further processing and ingestion, along with optional metadata.

Parameters
text (str): text to convert.
metadata (dict, optional): A dictionary containing metadata information associated with the text content. If not provided, 
a default metadata dictionary will be created as "text"

Return Value
The function returns a list of raw documents, where each document represents the plain text . Each document is represented as an instance of the Document class.
    """
    if metadata==None:
        metadata={"source": "text"} 
    documents = [text]
    raw_documents = [Document(page_content=raw_doc,metadata=metadata) for raw_doc in documents]
    return raw_documents
   

def convert_youtube_url(url):
    """
    The convert_youtube_url function is a Python function that converts a YouTube URL containing additional query parameters 
    to a simplified URL that only includes the video ID as a query parameter. 
    This can be useful for obtaining a consistent and concise representation of YouTube video URLs before ingestion.

Parameters
url (str): The YouTube URL that you want to convert.

Return Value
The function returns a new simplified YouTube URL that includes only the video ID as a query parameter.
    """
    parsed_url = urlparse(url)
    query_params = parse_qs(parsed_url.query)
    # Extract the video ID from the query parameters
    video_id = query_params.get('v', [''])[0]
    # Construct the new URL with the video ID
    new_query_params = {'v': video_id}
    new_parsed_url = parsed_url._replace(query=urlencode(new_query_params))
    new_url = urlunparse(new_parsed_url)
    return new_url

def convert_short_youtube_url(url):
    """
    The convert_short_youtube_url function is a Python function that converts a short YouTube URL (containing only the video ID in the path)
    to a full YouTube URL with additional query parameters. This can be useful for obtaining a more comprehensive and usable YouTube URL.

Parameters
url (str): The short YouTube URL that you want to convert.

Return Value
The function returns a new full YouTube URL that includes the video ID as a query parameter.
    """
    parsed_url = urlparse(url)

    # Extract the video ID from the path
    video_id = parsed_url.path.lstrip('/')

    # Construct the new URL with the video ID
    new_netloc = "www.youtube.com"
    new_path = "/watch"
    new_query_params = {'v': video_id}
    new_parsed_url = parsed_url._replace(netloc=new_netloc, path=new_path, query=urlencode(new_query_params))
    new_url = urlunparse(new_parsed_url)

    return new_url



def youtube_to_rawdocs(videoUrl, metadata=None):
    """
    The youtube_to_rawdocs function is a Python function that extracts the transcript from a YouTube video and converts it into a list of raw documents. 
    The function uses the youtube-transcript-api and pytube libraries to fetch the transcript from the video URL 
    and create raw documents with the extracted text.

Parameters
videoUrl (str): The URL of the YouTube video from which you want to extract the transcript.
metadata (dict, optional): A dictionary containing metadata information associated with the YouTube video. If not provided, a default metadata dictionary will be created with the video source information.

Return Value
The function returns a list of raw documents, where each document represents a segment of the video's transcript. Each document is represented as an instance of the Document class.
    """
    
    # pip install youtube-transcript-api
    # pip install pytube
    
    #example of video url: "https://www.youtube.com/watch?v=QsYGlZkevEg"
    # https://www.youtube.com/live/GYkq9Rgoj8E?feature=share to https://www.youtube.com/watch?v=GYkq9Rgoj8E
    if ("feature=share" in videoUrl):
        videoUrl=convert_youtube_url(videoUrl)
     #https://youtu.be/Lz28P5rbEzY to https://www.youtube.com/watch?v=Lz28P5rbEzY  
    if ("youtu.be" in videoUrl):
        videoUrl=convert_short_youtube_url(videoUrl)
   
    try:
        loader = YoutubeLoader.from_youtube_channel(videoUrl, add_video_info=False)
        raw_docs = loader.load()
        raw_documents = [Document(page_content=raw_doc.page_content,metadata=metadata) for raw_doc in raw_docs]
        return raw_documents
    except Exception as e :
        raise HTTPException(status_code=400, detail="Error with Youtube video " +videoUrl+": "+ str(e))   

def mp3_to_rawdocs(audioUrl, metadata=None):
    aai.settings.api_key ="0a1986a0bda24904bbcb538d2c5f60b5"
    transcriber = aai.Transcriber()
    transcript = transcriber.transcribe(audioUrl)
    raw_documents=string_to_rawdocs(transcript.text,None)
    # loader = AssemblyAIAudioTranscriptLoader(audioUrl,api_key="0a1986a0bda24904bbcb538d2c5f60b5")
    # raw_docs = loader.load()
    # if metadata==None:
    #         metadata={"source": audioUrl}
    # raw_documents = [Document(page_content=raw_doc.page_content,metadata=raw_doc.metadata) for raw_doc in raw_docs]
    return raw_documents

def verify_filename_before_ingestion(doc): 
    
        filename=doc.strip()
        filename=filename.replace('https:https:', 'https:')
        extension = os.path.splitext(filename)[1].lower()
        if extension==".pdf" or filename[-3:].lower()=="pdf":
            return True
        elif extension==".csv" or filename[-3:].lower()=="csv":
            return True
        elif extension=='.ppt' or extension=='.pptx' or filename[-3:].lower()=="ppt" or filename[-4:].lower()=="pptx":
            return True
        elif extension=='.doc' or extension=='.docx' or filename[-3:].lower()=="doc" or filename[-4:].lower()=="docx":
            return True
        elif extension=='.html' or extension=='.htm' or filename[-4:].lower()=="html" or filename[-3:].lower()=="htm":
            return True
        elif extension=='.txt'or filename[-3:].lower()=="txt":
            return True
        elif is_youtube_video(filename):
             return True
        else:
            return False   
        
def pinecone_namespace_to_vectorestore(pinecone_apik, open_apik, index_name,pinecone_env, ns)-> VectorStore:
   """
 The pinecone_namespace_to_vectorestore function is a Python function that retrieves a namespace from an existing Pinecone index, 
 converts it to a VectorStore object, and returns it. It uses Pinecone's Python SDK to interact with the index and OpenAI's embeddings 
 to manage the embeddings for the index.

Parameters
pinecone_apik (str): The API key for Pinecone, used for authentication and access to the Pinecone service.
open_apik (str): The API key for OpenAI, used for embedding retrieval.
index_name (str): The name of the Pinecone index to retrieve.
ns (str): The namespace associated with the index.

Return Value
Returns a VectorStore object representing the Pinecone index with the specified name and namespace.
   """
   try:
    pinecone.init(
        api_key=pinecone_apik,  # find at app.pinecone.io
        environment=pinecone_env,  
    )
    embeddings = OpenAIEmbeddings(openai_api_key=open_apik)
    return Pinecone.from_existing_index(index_name=index_name,embedding=embeddings,namespace=ns)
   except Exception:
       return None
   
def verify_filenames_before_ingestion(docs): 
    """
    The verify_filenames_before_ingestion function is a Python function that performs verification checks on a list of documents 
    before they are ingested. It ensures that the documents' URLs or filenames are valid and supported for ingestion based on their format. 
    If a document's URL or filename is not allowed or has an unsupported format, the function raises an exception.

Parameters:
docs (list of Documents urls): A list of Documents urls representing the documents to be ingested.

Return Value
This function does not have a return value. It raises an HTTPException if a document's URL or filename is not allowed or has an unsupported format.
    """
    for doc in docs:
        filename=doc.strip()
        filename=filename.replace('https:https:', 'https:')
        extension = os.path.splitext(filename)[1].lower()
        if extension==".pdf" or filename[-3:].lower()=="pdf":
            continue
        elif extension==".csv" or filename[-3:].lower()=="csv":
            continue
        elif extension=='.ppt' or extension=='.pptx' or filename[-3:].lower()=="ppt" or filename[-4:].lower()=="pptx":
            continue
        elif extension=='.doc' or extension=='.docx' or filename[-3:].lower()=="doc" or filename[-4:].lower()=="docx":
            continue
        elif extension=='.html' or extension=='.htm' or filename[-4:].lower()=="html" or filename[-3:].lower()=="htm":
            continue
        elif extension=='.txt'or filename[-3:].lower()=="txt":
            continue
        elif extension=='.mp3'or filename[-3:].lower()=="mp3":
            continue
        elif is_youtube_video(filename):
                    pass
        else:
            raise HTTPException(status_code=400, detail="File path \'" +filename+"\' is not a valid file or url, or format not yet supported")